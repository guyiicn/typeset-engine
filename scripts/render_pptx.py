#!/usr/bin/env python3
"""
PPTX 渲染引擎 — 支持多种投行/券商风格模板。

模板:
  default        — 简洁商务蓝
  cicc           — 中金红
  goldman        — 高盛蓝灰
  morgan         — 摩根深蓝
  dark           — 深色主题
  minimal        — 极简白

数据格式 (JSON):
{
  "title": "贵州茅台 研究报告",
  "subtitle": "600519 | 白酒/消费",
  "author": "FinRobot",
  "date": "2026-04-07",
  "slides": [
    {
      "layout": "title",          # title / section / content / two_column / chart / table / summary
      "title": "公司概况",
      "content": "正文内容...",
      "bullets": ["要点1", "要点2"],
      "image": "/path/to/chart.png",
      "table": {"headers": [...], "rows": [...]},
      "notes": "演讲者备注"
    }
  ]
}
"""

import os
from pathlib import Path
from typing import Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ═══════════════════════════════════════════
# 主题配置
# ═══════════════════════════════════════════

THEMES = {
    'default': {
        'name': 'Business Blue',
        'primary': RGBColor(0x1a, 0x36, 0x5d),      # 深蓝
        'secondary': RGBColor(0xc9, 0xa2, 0x27),     # 金色
        'text_dark': RGBColor(0x2d, 0x37, 0x48),
        'text_light': RGBColor(0xff, 0xff, 0xff),
        'bg_color': RGBColor(0xff, 0xff, 0xff),
        'accent_bg': RGBColor(0x1a, 0x36, 0x5d),
        'table_header': RGBColor(0x1a, 0x36, 0x5d),
        'table_alt': RGBColor(0xf0, 0xf4, 0xf8),
        'font_title': 'Noto Sans CJK SC',
        'font_body': 'Noto Sans CJK SC',
        'font_title_fallback': 'Arial',
        'font_body_fallback': 'Arial',
    },
    'cicc': {
        'name': '中金研究',
        'primary': RGBColor(0xC4, 0x1E, 0x3A),      # 中金红
        'secondary': RGBColor(0x8B, 0x00, 0x00),     # 深红
        'text_dark': RGBColor(0x33, 0x33, 0x33),
        'text_light': RGBColor(0xff, 0xff, 0xff),
        'bg_color': RGBColor(0xff, 0xff, 0xff),
        'accent_bg': RGBColor(0xC4, 0x1E, 0x3A),
        'table_header': RGBColor(0xf5, 0xf5, 0xf5),
        'table_alt': RGBColor(0xfa, 0xfa, 0xfa),
        'font_title': 'Noto Sans CJK SC',
        'font_body': 'Noto Sans CJK SC',
        'font_title_fallback': 'Microsoft YaHei',
        'font_body_fallback': 'Microsoft YaHei',
    },
    'goldman': {
        'name': 'Goldman Sachs',
        'primary': RGBColor(0x00, 0x3A, 0x70),      # GS深蓝
        'secondary': RGBColor(0x6B, 0x8E, 0x23),     # 橄榄绿
        'text_dark': RGBColor(0x33, 0x33, 0x33),
        'text_light': RGBColor(0xff, 0xff, 0xff),
        'bg_color': RGBColor(0xff, 0xff, 0xff),
        'accent_bg': RGBColor(0x00, 0x3A, 0x70),
        'table_header': RGBColor(0x00, 0x3A, 0x70),
        'table_alt': RGBColor(0xf0, 0xf5, 0xfa),
        'font_title': 'Arial',
        'font_body': 'Arial',
        'font_title_fallback': 'Helvetica',
        'font_body_fallback': 'Helvetica',
    },
    'morgan': {
        'name': 'Morgan Stanley',
        'primary': RGBColor(0x00, 0x1E, 0x62),      # MS深蓝
        'secondary': RGBColor(0x00, 0x7A, 0xB8),     # 亮蓝
        'text_dark': RGBColor(0x33, 0x33, 0x33),
        'text_light': RGBColor(0xff, 0xff, 0xff),
        'bg_color': RGBColor(0xff, 0xff, 0xff),
        'accent_bg': RGBColor(0x00, 0x1E, 0x62),
        'table_header': RGBColor(0x00, 0x1E, 0x62),
        'table_alt': RGBColor(0xf0, 0xf0, 0xf8),
        'font_title': 'Arial',
        'font_body': 'Arial',
        'font_title_fallback': 'Helvetica',
        'font_body_fallback': 'Helvetica',
    },
    'dark': {
        'name': 'Dark Theme',
        'primary': RGBColor(0x00, 0xd4, 0xaa),      # 青绿
        'secondary': RGBColor(0xff, 0x6b, 0x6b),     # 珊瑚红
        'text_dark': RGBColor(0xe0, 0xe0, 0xe0),
        'text_light': RGBColor(0xff, 0xff, 0xff),
        'bg_color': RGBColor(0x1a, 0x1a, 0x2e),
        'accent_bg': RGBColor(0x16, 0x21, 0x3e),
        'table_header': RGBColor(0x16, 0x21, 0x3e),
        'table_alt': RGBColor(0x20, 0x2a, 0x44),
        'font_title': 'Noto Sans CJK SC',
        'font_body': 'Noto Sans CJK SC',
        'font_title_fallback': 'Arial',
        'font_body_fallback': 'Arial',
    },
    'minimal': {
        'name': 'Minimal White',
        'primary': RGBColor(0x33, 0x33, 0x33),
        'secondary': RGBColor(0x99, 0x99, 0x99),
        'text_dark': RGBColor(0x33, 0x33, 0x33),
        'text_light': RGBColor(0xff, 0xff, 0xff),
        'bg_color': RGBColor(0xff, 0xff, 0xff),
        'accent_bg': RGBColor(0x33, 0x33, 0x33),
        'table_header': RGBColor(0xf5, 0xf5, 0xf5),
        'table_alt': RGBColor(0xfa, 0xfa, 0xfa),
        'font_title': 'Noto Sans CJK SC',
        'font_body': 'Noto Sans CJK SC',
        'font_title_fallback': 'Arial',
        'font_body_fallback': 'Arial',
    },
}


# ═══════════════════════════════════════════
# PPT 生成核心
# ═══════════════════════════════════════════

class PptxBuilder:
    """PPTX 文档构建器"""

    def __init__(self, theme_name: str = 'default'):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)   # 16:9
        self.prs.slide_height = Inches(7.5)
        self.theme = THEMES.get(theme_name, THEMES['default'])
        self.slide_w = self.prs.slide_width
        self.slide_h = self.prs.slide_height

    def _add_bg(self, slide, color=None):
        """设置幻灯片背景色"""
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color or self.theme['bg_color']

    def _add_text(self, slide, text, left, top, width, height,
                   font_size=18, color=None, bold=False, alignment=PP_ALIGN.LEFT,
                   font_name=None):
        """添加文本框"""
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color or self.theme['text_dark']
        p.font.bold = bold
        p.font.name = font_name or self.theme['font_title']
        p.alignment = alignment
        return txBox

    def _add_shape(self, slide, shape_type, left, top, width, height, fill_color):
        """添加形状"""
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        return shape

    # ── 幻灯片布局 ──

    def add_title_slide(self, title: str, subtitle: str = '', author: str = '', date: str = ''):
        """封面页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # blank
        self._add_bg(slide)

        # 顶部色条
        self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0), self.slide_w, Inches(0.15),
                        self.theme['primary'])

        # 底部色块
        self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(4.5), self.slide_w, Inches(3),
                        self.theme['accent_bg'])

        # 标题
        self._add_text(slide, title,
                        Inches(1), Inches(1.5), Inches(11), Inches(1.5),
                        font_size=40, bold=True, color=self.theme['primary'],
                        alignment=PP_ALIGN.LEFT)

        # 副标题
        if subtitle:
            self._add_text(slide, subtitle,
                            Inches(1), Inches(3.2), Inches(11), Inches(0.8),
                            font_size=20, color=self.theme['secondary'])

        # 作者和日期
        info = f"{author}  |  {date}" if author else date
        if info:
            self._add_text(slide, info,
                            Inches(1), Inches(5.5), Inches(11), Inches(0.6),
                            font_size=16, color=self.theme['text_light'],
                            alignment=PP_ALIGN.LEFT)

    def add_section_slide(self, title: str, subtitle: str = ''):
        """章节分隔页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, self.theme['accent_bg'])

        # 左侧竖条
        self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                        Inches(0.8), Inches(2), Inches(0.08), Inches(2),
                        self.theme['text_light'])

        self._add_text(slide, title,
                        Inches(1.2), Inches(2.5), Inches(10), Inches(1.2),
                        font_size=36, bold=True, color=self.theme['text_light'])

        if subtitle:
            self._add_text(slide, subtitle,
                            Inches(1.2), Inches(4), Inches(10), Inches(0.8),
                            font_size=18, color=self.theme['secondary'])

    def add_content_slide(self, title: str, content: str = '',
                           bullets: List[str] = None, image: str = None,
                           notes: str = None):
        """正文页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)

        # 顶部色条
        self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0), self.slide_w, Inches(0.06),
                        self.theme['primary'])

        # 标题
        self._add_text(slide, title,
                        Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                        font_size=28, bold=True, color=self.theme['primary'])

        # 标题下分隔线
        self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                        Inches(0.8), Inches(1.15), Inches(11.5), Inches(0.02),
                        self.theme['primary'])

        content_top = Inches(1.5)
        content_width = Inches(7.5) if image else Inches(11.5)

        # 正文
        if content:
            self._add_text(slide, content,
                            Inches(0.8), content_top, content_width, Inches(4.5),
                            font_size=16, color=self.theme['text_dark'],
                            font_name=self.theme['font_body'])

        # 要点列表
        if bullets:
            txBox = slide.shapes.add_textbox(
                Inches(0.8), content_top + (Inches(1.5) if content else Inches(0)),
                content_width, Inches(4))
            tf = txBox.text_frame
            tf.word_wrap = True
            for i, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"  •  {bullet}"
                p.font.size = Pt(15)
                p.font.color.rgb = self.theme['text_dark']
                p.font.name = self.theme['font_body']
                p.space_after = Pt(8)

        # 图片
        if image and os.path.exists(image):
            try:
                slide.shapes.add_picture(image,
                    Inches(8.8), Inches(1.5), Inches(4), Inches(3))
            except:
                pass

        # 演讲者备注
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def add_two_column_slide(self, title: str, left_content: str,
                              right_content: str = '', right_image: str = None):
        """双栏页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)

        self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0), self.slide_w, Inches(0.06),
                        self.theme['primary'])

        self._add_text(slide, title,
                        Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                        font_size=28, bold=True, color=self.theme['primary'])

        self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                        Inches(0.8), Inches(1.15), Inches(11.5), Inches(0.02),
                        self.theme['primary'])

        # 左栏
        self._add_text(slide, left_content,
                        Inches(0.8), Inches(1.5), Inches(5.5), Inches(5),
                        font_size=14, color=self.theme['text_dark'],
                        font_name=self.theme['font_body'])

        # 右栏
        if right_image and os.path.exists(right_image):
            slide.shapes.add_picture(right_image,
                Inches(7), Inches(1.5), Inches(5.5), Inches(4))
        elif right_content:
            self._add_text(slide, right_content,
                            Inches(7), Inches(1.5), Inches(5.5), Inches(5),
                            font_size=14, color=self.theme['text_dark'],
                            font_name=self.theme['font_body'])

    def add_table_slide(self, title: str, headers: List[str], rows: List[List[str]]):
        """表格页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)

        self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0), self.slide_w, Inches(0.06),
                        self.theme['primary'])

        self._add_text(slide, title,
                        Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                        font_size=28, bold=True, color=self.theme['primary'])

        # 表格
        n_rows = min(len(rows) + 1, 15)  # 最多 15 行
        n_cols = len(headers)
        table_shape = slide.shapes.add_table(
            n_rows, n_cols,
            Inches(0.8), Inches(1.5),
            Inches(11.5), Inches(5)
        )
        table = table_shape.table

        # 表头
        for j, h in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = str(h)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.theme['table_header']
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.bold = True
                p.font.color.rgb = self.theme['text_light'] if self.theme['table_header'] != self.theme['table_alt'] else self.theme['text_dark']
                p.font.name = self.theme['font_body']

        # 数据行
        for i, row in enumerate(rows[:14]):
            for j, val in enumerate(row):
                cell = table.cell(i + 1, j)
                cell.text = str(val)
                if i % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.theme['table_alt']
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(11)
                    p.font.color.rgb = self.theme['text_dark']
                    p.font.name = self.theme['font_body']

    def add_summary_slide(self, title: str, points: List[str]):
        """总结页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, self.theme['accent_bg'])

        self._add_text(slide, title,
                        Inches(1), Inches(0.8), Inches(11), Inches(1),
                        font_size=32, bold=True, color=self.theme['text_light'],
                        alignment=PP_ALIGN.CENTER)

        y = Inches(2.2)
        for point in points:
            self._add_text(slide, f"✓  {point}",
                            Inches(2), y, Inches(9), Inches(0.6),
                            font_size=18, color=self.theme['text_light'])
            y += Inches(0.7)

    def add_kpi_slide(self, title: str, kpis: List[Dict]):
        """
        指标卡片页 — 4-6个 KPI 大数字展示。
        kpis = [{"label": "营收", "value": "¥1,709亿", "change": "+15.7%"}, ...]
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)

        self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0), self.slide_w, Inches(0.06),
                        self.theme['primary'])

        self._add_text(slide, title,
                        Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                        font_size=28, bold=True, color=self.theme['primary'])

        self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                        Inches(0.8), Inches(1.15), Inches(11.5), Inches(0.02),
                        self.theme['primary'])

        # 布局 KPI 卡片（每行3个）
        cols = min(len(kpis), 3)
        card_w = 11.5 / cols
        for i, kpi in enumerate(kpis[:6]):
            col = i % cols
            row = i // cols
            x = Inches(0.8 + col * card_w)
            y = Inches(1.8 + row * 2.5)

            # 卡片背景
            self._add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                            x, y, Inches(card_w - 0.3), Inches(2),
                            self.theme['table_alt'] if self.theme['bg_color'] == RGBColor(0xff, 0xff, 0xff)
                            else self.theme['accent_bg'])

            # 数值
            self._add_text(slide, kpi.get('value', ''),
                            x + Inches(0.2), y + Inches(0.3),
                            Inches(card_w - 0.7), Inches(0.8),
                            font_size=32, bold=True, color=self.theme['primary'],
                            alignment=PP_ALIGN.CENTER)

            # 标签
            self._add_text(slide, kpi.get('label', ''),
                            x + Inches(0.2), y + Inches(1.1),
                            Inches(card_w - 0.7), Inches(0.4),
                            font_size=14, color=self.theme['text_dark'],
                            alignment=PP_ALIGN.CENTER,
                            font_name=self.theme['font_body'])

            # 变化率
            change = kpi.get('change', '')
            if change:
                is_positive = change.startswith('+')
                change_color = RGBColor(0xCC, 0x00, 0x00) if is_positive else RGBColor(0x00, 0x99, 0x00)
                self._add_text(slide, change,
                                x + Inches(0.2), y + Inches(1.5),
                                Inches(card_w - 0.7), Inches(0.3),
                                font_size=12, color=change_color,
                                alignment=PP_ALIGN.CENTER)

    def add_chart_slide(self, title: str, image: str, caption: str = '',
                         notes_text: str = ''):
        """全幅图表页 — 一张图占满幻灯片"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)

        self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0), self.slide_w, Inches(0.06),
                        self.theme['primary'])

        self._add_text(slide, title,
                        Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                        font_size=28, bold=True, color=self.theme['primary'])

        if image and os.path.exists(image):
            slide.shapes.add_picture(image,
                Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.2))

        if caption:
            self._add_text(slide, caption,
                            Inches(0.8), Inches(6.7), Inches(11), Inches(0.4),
                            font_size=10, color=self.theme['secondary'],
                            font_name=self.theme['font_body'])

        if notes_text:
            slide.notes_slide.notes_text_frame.text = notes_text

    def add_comparison_slide(self, title: str, left_title: str, left_items: List[str],
                              right_title: str, right_items: List[str]):
        """对比页 — 左右两栏对比（优势vs风险、现在vs未来等）"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)

        self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0), self.slide_w, Inches(0.06),
                        self.theme['primary'])

        self._add_text(slide, title,
                        Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                        font_size=28, bold=True, color=self.theme['primary'])

        self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                        Inches(0.8), Inches(1.15), Inches(11.5), Inches(0.02),
                        self.theme['primary'])

        # 中间分隔线
        self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                        Inches(6.6), Inches(1.5), Inches(0.03), Inches(5),
                        self.theme['secondary'])

        # 左栏标题
        self._add_text(slide, left_title,
                        Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.6),
                        font_size=20, bold=True, color=self.theme['primary'])

        # 左栏内容
        y = Inches(2.3)
        for item in left_items[:8]:
            self._add_text(slide, f"▸ {item}",
                            Inches(1), y, Inches(5), Inches(0.5),
                            font_size=14, color=self.theme['text_dark'],
                            font_name=self.theme['font_body'])
            y += Inches(0.5)

        # 右栏标题
        self._add_text(slide, right_title,
                        Inches(7), Inches(1.5), Inches(5.5), Inches(0.6),
                        font_size=20, bold=True, color=self.theme['secondary'])

        # 右栏内容
        y = Inches(2.3)
        for item in right_items[:8]:
            self._add_text(slide, f"▸ {item}",
                            Inches(7.2), y, Inches(5), Inches(0.5),
                            font_size=14, color=self.theme['text_dark'],
                            font_name=self.theme['font_body'])
            y += Inches(0.5)

    def add_timeline_slide(self, title: str, events: List[Dict]):
        """
        时间线页 — 展示里程碑/事件序列。
        events = [{"date": "2024 Q1", "event": "营收突破1500亿"}, ...]
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)

        self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0), self.slide_w, Inches(0.06),
                        self.theme['primary'])

        self._add_text(slide, title,
                        Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                        font_size=28, bold=True, color=self.theme['primary'])

        # 时间线横轴
        n = min(len(events), 6)
        if n == 0:
            return

        line_y = Inches(4)
        self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                        Inches(1), line_y, Inches(11.3), Inches(0.03),
                        self.theme['primary'])

        step = 11.3 / n
        for i, evt in enumerate(events[:6]):
            x = Inches(1 + i * step + step / 2 - 0.5)

            # 节点圆点
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, x + Inches(0.4), line_y - Inches(0.08),
                Inches(0.2), Inches(0.2))
            dot.fill.solid()
            dot.fill.fore_color.rgb = self.theme['primary']
            dot.line.fill.background()

            # 日期（上方）
            self._add_text(slide, evt.get('date', ''),
                            x, line_y - Inches(0.8), Inches(1), Inches(0.5),
                            font_size=11, bold=True, color=self.theme['primary'],
                            alignment=PP_ALIGN.CENTER)

            # 事件（下方）
            self._add_text(slide, evt.get('event', ''),
                            x - Inches(0.2), line_y + Inches(0.3), Inches(1.4), Inches(1.5),
                            font_size=10, color=self.theme['text_dark'],
                            alignment=PP_ALIGN.CENTER,
                            font_name=self.theme['font_body'])

    def add_quote_slide(self, quote: str, author: str = '', source: str = ''):
        """引用页 — 大字引用"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, self.theme['accent_bg'])

        # 引号装饰
        self._add_text(slide, '"',
                        Inches(1), Inches(1), Inches(2), Inches(1.5),
                        font_size=80, color=self.theme['secondary'],
                        alignment=PP_ALIGN.LEFT)

        # 引用文本
        self._add_text(slide, quote,
                        Inches(1.5), Inches(2.5), Inches(10), Inches(2.5),
                        font_size=24, color=self.theme['text_light'],
                        alignment=PP_ALIGN.LEFT)

        # 作者
        attribution = f"— {author}" + (f", {source}" if source else "")
        if author:
            self._add_text(slide, attribution,
                            Inches(1.5), Inches(5.5), Inches(10), Inches(0.6),
                            font_size=16, color=self.theme['secondary'],
                            alignment=PP_ALIGN.RIGHT)

    def add_end_slide(self, title: str = 'Thank You', subtitle: str = '',
                       contact: str = ''):
        """结束页"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, self.theme['accent_bg'])

        self._add_text(slide, title,
                        Inches(1), Inches(2.5), Inches(11), Inches(1.5),
                        font_size=44, bold=True, color=self.theme['text_light'],
                        alignment=PP_ALIGN.CENTER)

        if subtitle:
            self._add_text(slide, subtitle,
                            Inches(1), Inches(4.2), Inches(11), Inches(0.8),
                            font_size=18, color=self.theme['secondary'],
                            alignment=PP_ALIGN.CENTER)

        if contact:
            self._add_text(slide, contact,
                            Inches(1), Inches(5.5), Inches(11), Inches(0.6),
                            font_size=14, color=self.theme['text_light'],
                            alignment=PP_ALIGN.CENTER)

    # ═══════════════════════════════════════════
    # 投行 Pitch Book 专用 Slide 类型
    # ═══════════════════════════════════════════

    def add_comparable_companies_slide(self, title: str, headers: List[str],
                                        rows: List[List[str]],
                                        summary_rows: List[Dict] = None,
                                        source: str = ''):
        """可比公司分析 — 投行核心 slide，宽表格 10-15 列，带 Median/Mean 汇总行"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)

        # 顶部色条
        self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0), self.slide_w, Inches(0.06),
                        self.theme['primary'])

        # 标题
        self._add_text(slide, title,
                        Inches(0.5), Inches(0.25), Inches(12), Inches(0.5),
                        font_size=22, bold=True, color=self.theme['primary'])

        # 标题下划线
        self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                        Inches(0.5), Inches(0.8), Inches(12.2), Inches(0.015),
                        self.theme['primary'])

        # 表格 — 紧凑布局，字号 8-9pt
        n_summary = len(summary_rows) if summary_rows else 0
        n_rows = min(len(rows) + 1 + n_summary, 22)  # 表头 + 数据 + 汇总
        n_cols = len(headers)

        table_shape = slide.shapes.add_table(
            n_rows, n_cols,
            Inches(0.3), Inches(1.0),
            Inches(12.6), Inches(5.5)
        )
        table = table_shape.table

        # 表头
        for j, h in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = str(h)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.theme['table_header']
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(8)
                p.font.bold = True
                p.font.color.rgb = self.theme['text_light']
                p.font.name = self.theme['font_body_fallback']
                p.alignment = PP_ALIGN.CENTER

        # 数据行
        max_data = n_rows - 1 - n_summary
        for i, row in enumerate(rows[:max_data]):
            for j, val in enumerate(row):
                cell = table.cell(i + 1, j)
                cell.text = str(val)
                if i % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.theme['table_alt']
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(8)
                    p.font.color.rgb = self.theme['text_dark']
                    p.font.name = self.theme['font_body_fallback']
                    # 第一列左对齐（公司名），其余居中
                    p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER

        # 汇总行（Median/Mean）— 加粗 + 顶部线
        if summary_rows:
            data_end = len(rows[:max_data]) + 1
            for si, sr in enumerate(summary_rows):
                row_idx = data_end + si
                if row_idx >= n_rows:
                    break
                values = [sr.get('label', '')] + sr.get('values', [])
                for j, val in enumerate(values[:n_cols]):
                    cell = table.cell(row_idx, j)
                    cell.text = str(val)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.theme['table_alt']
                    for p in cell.text_frame.paragraphs:
                        p.font.size = Pt(8)
                        p.font.bold = True
                        p.font.color.rgb = self.theme['primary']
                        p.font.name = self.theme['font_body_fallback']
                        p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER

        # Source 脚注
        if source:
            self._add_text(slide, f"Source: {source}",
                            Inches(0.5), Inches(6.7), Inches(12), Inches(0.3),
                            font_size=7, color=RGBColor(0x80, 0x80, 0x80),
                            font_name=self.theme['font_body_fallback'])

    def add_football_field_slide(self, title: str, ranges: List[Dict],
                                   current_price: float = None,
                                   currency: str = '$', source: str = ''):
        """估值区间图（Football Field）— 横向条形图，多种估值方法对比

        ranges = [
            {"method": "52-Week Range", "low": 28.5, "high": 45.2},
            {"method": "DCF Analysis", "low": 35.0, "high": 50.0},
        ]
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)

        self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0), self.slide_w, Inches(0.06),
                        self.theme['primary'])

        self._add_text(slide, title,
                        Inches(0.8), Inches(0.25), Inches(11), Inches(0.5),
                        font_size=22, bold=True, color=self.theme['primary'])

        self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                        Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.015),
                        self.theme['primary'])

        if not ranges:
            return

        # 计算全局范围
        all_vals = [r['low'] for r in ranges] + [r['high'] for r in ranges]
        if current_price is not None:
            all_vals.append(current_price)
        global_min = min(all_vals) * 0.85
        global_max = max(all_vals) * 1.1
        scale = 9.5 / (global_max - global_min)  # 图表区域 9.5 inches

        n = len(ranges)
        bar_h = min(0.5, 4.0 / n)  # 每条高度
        gap = 0.15
        start_y = 1.5
        chart_left = 3.0  # 留空给方法名标签

        # 颜色梯度
        bar_colors = [
            self.theme['primary'],
            self.theme['secondary'],
            RGBColor(0x54, 0x82, 0x35),  # 绿
            RGBColor(0x80, 0x80, 0x80),  # 灰
            RGBColor(0xE5, 0x3E, 0x3E),  # 红
            RGBColor(0x7C, 0x3A, 0xED),  # 紫
        ]

        for i, r in enumerate(ranges):
            y = start_y + i * (bar_h + gap)
            low, high = r['low'], r['high']

            # 方法名标签
            self._add_text(slide, r['method'],
                            Inches(0.5), Inches(y), Inches(2.3), Inches(bar_h),
                            font_size=10, color=self.theme['text_dark'],
                            font_name=self.theme['font_body_fallback'],
                            alignment=PP_ALIGN.RIGHT)

            # 条形
            bar_left = chart_left + (low - global_min) * scale
            bar_width = (high - low) * scale
            color = bar_colors[i % len(bar_colors)]

            bar = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(bar_left), Inches(y + 0.05),
                Inches(max(bar_width, 0.1)), Inches(bar_h - 0.1))
            bar.fill.solid()
            bar.fill.fore_color.rgb = color
            bar.line.fill.background()

            # 数值标注（低-高）
            self._add_text(slide, f"{currency}{low:.1f}",
                            Inches(bar_left - 0.6), Inches(y), Inches(0.55), Inches(bar_h),
                            font_size=8, color=self.theme['text_dark'],
                            alignment=PP_ALIGN.RIGHT,
                            font_name=self.theme['font_body_fallback'])
            self._add_text(slide, f"{currency}{high:.1f}",
                            Inches(bar_left + bar_width + 0.05), Inches(y),
                            Inches(0.55), Inches(bar_h),
                            font_size=8, color=self.theme['text_dark'],
                            alignment=PP_ALIGN.LEFT,
                            font_name=self.theme['font_body_fallback'])

        # 当前价格竖线
        if current_price is not None:
            line_x = chart_left + (current_price - global_min) * scale
            total_h = n * (bar_h + gap)
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(line_x), Inches(start_y - 0.1),
                Inches(0.02), Inches(total_h + 0.3))
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(0xE5, 0x3E, 0x3E)
            line.line.fill.background()

            self._add_text(slide, f"Current: {currency}{current_price:.1f}",
                            Inches(line_x - 0.5), Inches(start_y + total_h + 0.2),
                            Inches(1.2), Inches(0.3),
                            font_size=9, bold=True, color=RGBColor(0xE5, 0x3E, 0x3E),
                            alignment=PP_ALIGN.CENTER,
                            font_name=self.theme['font_body_fallback'])

        # Source
        if source:
            self._add_text(slide, f"Source: {source}",
                            Inches(0.5), Inches(6.7), Inches(12), Inches(0.3),
                            font_size=7, color=RGBColor(0x80, 0x80, 0x80),
                            font_name=self.theme['font_body_fallback'])

    def add_sources_uses_slide(self, title: str, sources: List[Dict],
                                 uses: List[Dict], currency: str = '$m',
                                 source_note: str = ''):
        """资金来源与用途（Sources & Uses）— 双栏对称表格

        sources = [{"item": "Term Loan A", "amount": 500}, ...]
        uses = [{"item": "Purchase Equity", "amount": 850}, ...]
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)

        self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0), self.slide_w, Inches(0.06),
                        self.theme['primary'])

        self._add_text(slide, title,
                        Inches(0.8), Inches(0.25), Inches(11), Inches(0.5),
                        font_size=22, bold=True, color=self.theme['primary'])

        self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                        Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.015),
                        self.theme['primary'])

        # Sources 表格（左）
        s_total = sum(s['amount'] for s in sources)
        s_rows = len(sources) + 2  # header + data + total
        s_table = slide.shapes.add_table(
            s_rows, 2, Inches(0.8), Inches(1.2), Inches(5.5), Inches(4.5)
        ).table

        # Sources 表头
        for j, h in enumerate([f'Sources', f'Amount ({currency})']):
            cell = s_table.cell(0, j)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.theme['primary']
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.bold = True
                p.font.color.rgb = self.theme['text_light']
                p.font.name = self.theme['font_body_fallback']

        # Sources 数据
        for i, s in enumerate(sources):
            s_table.cell(i+1, 0).text = s['item']
            s_table.cell(i+1, 1).text = f"{s['amount']:,.0f}"
            for j in range(2):
                for p in s_table.cell(i+1, j).text_frame.paragraphs:
                    p.font.size = Pt(10)
                    p.font.color.rgb = self.theme['text_dark']
                    p.font.name = self.theme['font_body_fallback']
                    p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT

        # Sources Total
        s_table.cell(s_rows-1, 0).text = 'Total Sources'
        s_table.cell(s_rows-1, 1).text = f"{s_total:,.0f}"
        for j in range(2):
            cell = s_table.cell(s_rows-1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.theme['table_alt']
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.bold = True
                p.font.color.rgb = self.theme['primary']
                p.font.name = self.theme['font_body_fallback']
                p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT

        # Uses 表格（右）
        u_total = sum(u['amount'] for u in uses)
        u_rows = len(uses) + 2
        u_table = slide.shapes.add_table(
            u_rows, 2, Inches(7), Inches(1.2), Inches(5.5), Inches(4.5)
        ).table

        for j, h in enumerate([f'Uses', f'Amount ({currency})']):
            cell = u_table.cell(0, j)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.theme['primary']
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.bold = True
                p.font.color.rgb = self.theme['text_light']
                p.font.name = self.theme['font_body_fallback']

        for i, u in enumerate(uses):
            u_table.cell(i+1, 0).text = u['item']
            u_table.cell(i+1, 1).text = f"{u['amount']:,.0f}"
            for j in range(2):
                for p in u_table.cell(i+1, j).text_frame.paragraphs:
                    p.font.size = Pt(10)
                    p.font.color.rgb = self.theme['text_dark']
                    p.font.name = self.theme['font_body_fallback']
                    p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT

        u_table.cell(u_rows-1, 0).text = 'Total Uses'
        u_table.cell(u_rows-1, 1).text = f"{u_total:,.0f}"
        for j in range(2):
            cell = u_table.cell(u_rows-1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.theme['table_alt']
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.bold = True
                p.font.color.rgb = self.theme['primary']
                p.font.name = self.theme['font_body_fallback']
                p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.RIGHT

        # Source note
        if source_note:
            self._add_text(slide, f"Source: {source_note}",
                            Inches(0.5), Inches(6.7), Inches(12), Inches(0.3),
                            font_size=7, color=RGBColor(0x80, 0x80, 0x80),
                            font_name=self.theme['font_body_fallback'])

    def add_sensitivity_matrix_slide(self, title: str, row_label: str,
                                        col_label: str, row_values: List,
                                        col_values: List, matrix: List[List],
                                        highlight_row: int = None,
                                        highlight_col: int = None,
                                        source: str = ''):
        """敏感性分析矩阵 — WACC × Terminal Growth Rate 二维表

        row_values = [1.0, 1.5, 2.0, 2.5]     # Terminal Growth Rate
        col_values = [8.0, 9.0, 10.0, 11.0]    # WACC
        matrix = [[42, 38, 35, 32], [45, 41, 37, 34], ...]
        highlight_row/col = base case 位置索引
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)

        self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0), self.slide_w, Inches(0.06),
                        self.theme['primary'])

        self._add_text(slide, title,
                        Inches(0.8), Inches(0.25), Inches(11), Inches(0.5),
                        font_size=22, bold=True, color=self.theme['primary'])

        self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                        Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.015),
                        self.theme['primary'])

        n_rows = len(row_values) + 1  # +1 for col header
        n_cols = len(col_values) + 1  # +1 for row header

        # 列标签
        self._add_text(slide, col_label,
                        Inches(4), Inches(1.2), Inches(6), Inches(0.4),
                        font_size=12, bold=True, color=self.theme['primary'],
                        alignment=PP_ALIGN.CENTER)

        # 行标签（竖向）
        self._add_text(slide, row_label,
                        Inches(1.5), Inches(3.5), Inches(1.5), Inches(0.4),
                        font_size=12, bold=True, color=self.theme['primary'],
                        alignment=PP_ALIGN.CENTER)

        # 表格
        table_shape = slide.shapes.add_table(
            n_rows, n_cols,
            Inches(3), Inches(1.7),
            Inches(8), Inches(4.5)
        )
        tbl = table_shape.table

        # 左上角空白
        cell = tbl.cell(0, 0)
        cell.text = f"{row_label} \\ {col_label}"
        cell.fill.solid()
        cell.fill.fore_color.rgb = self.theme['primary']
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(8)
            p.font.bold = True
            p.font.color.rgb = self.theme['text_light']
            p.font.name = self.theme['font_body_fallback']
            p.alignment = PP_ALIGN.CENTER

        # 列头
        for j, cv in enumerate(col_values):
            cell = tbl.cell(0, j + 1)
            cell.text = f"{cv}%"
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.theme['primary']
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.bold = True
                p.font.color.rgb = self.theme['text_light']
                p.font.name = self.theme['font_body_fallback']
                p.alignment = PP_ALIGN.CENTER

        # 行头 + 数据
        for i, rv in enumerate(row_values):
            # 行头
            cell = tbl.cell(i + 1, 0)
            cell.text = f"{rv}%"
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.theme['table_alt']
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.bold = True
                p.font.color.rgb = self.theme['primary']
                p.font.name = self.theme['font_body_fallback']
                p.alignment = PP_ALIGN.CENTER

            # 数据
            for j, val in enumerate(matrix[i] if i < len(matrix) else []):
                cell = tbl.cell(i + 1, j + 1)
                cell.text = f"${val}" if isinstance(val, (int, float)) else str(val)

                # Base case 高亮
                is_highlight = (highlight_row is not None and highlight_col is not None
                                and i == highlight_row and j == highlight_col)
                if is_highlight:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.theme['primary']
                    fg_color = self.theme['text_light']
                else:
                    fg_color = self.theme['text_dark']

                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(11)
                    p.font.color.rgb = fg_color
                    p.font.name = self.theme['font_body_fallback']
                    p.alignment = PP_ALIGN.CENTER
                    if is_highlight:
                        p.font.bold = True

        if source:
            self._add_text(slide, f"Source: {source}",
                            Inches(0.5), Inches(6.7), Inches(12), Inches(0.3),
                            font_size=7, color=RGBColor(0x80, 0x80, 0x80),
                            font_name=self.theme['font_body_fallback'])

    def add_transaction_overview_slide(self, title: str, key_points: List[str],
                                        terms: List[Dict], source: str = ''):
        """交易概览 — 左侧文字要点 + 右侧关键条款表格

        key_points = ["要点1", "要点2", ...]
        terms = [{"term": "Transaction Value", "value": "$1.5bn"}, ...]
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)

        self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0), self.slide_w, Inches(0.06),
                        self.theme['primary'])

        self._add_text(slide, title,
                        Inches(0.8), Inches(0.25), Inches(11), Inches(0.5),
                        font_size=22, bold=True, color=self.theme['primary'])

        self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                        Inches(0.8), Inches(0.8), Inches(11.5), Inches(0.015),
                        self.theme['primary'])

        # 左侧要点（40%）
        self._add_text(slide, "Key Highlights",
                        Inches(0.8), Inches(1.2), Inches(5), Inches(0.4),
                        font_size=14, bold=True, color=self.theme['primary'])

        y = Inches(1.8)
        for pt in key_points[:8]:
            self._add_text(slide, f"•  {pt}",
                            Inches(1.0), y, Inches(4.8), Inches(0.45),
                            font_size=11, color=self.theme['text_dark'],
                            font_name=self.theme['font_body'])
            y += Inches(0.5)

        # 右侧条款表格（55%）
        if terms:
            n_rows = len(terms) + 1
            t_shape = slide.shapes.add_table(
                n_rows, 2, Inches(6.5), Inches(1.2), Inches(6), Inches(5)
            )
            tbl = t_shape.table

            for j, h in enumerate(["Term", "Details"]):
                cell = tbl.cell(0, j)
                cell.text = h
                cell.fill.solid()
                cell.fill.fore_color.rgb = self.theme['primary']
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(10)
                    p.font.bold = True
                    p.font.color.rgb = self.theme['text_light']
                    p.font.name = self.theme['font_body_fallback']

            for i, t in enumerate(terms):
                tbl.cell(i+1, 0).text = t.get('term', '')
                tbl.cell(i+1, 1).text = t.get('value', '')
                for j in range(2):
                    if i % 2 == 1:
                        tbl.cell(i+1, j).fill.solid()
                        tbl.cell(i+1, j).fill.fore_color.rgb = self.theme['table_alt']
                    for p in tbl.cell(i+1, j).text_frame.paragraphs:
                        p.font.size = Pt(10)
                        p.font.color.rgb = self.theme['text_dark']
                        p.font.name = self.theme['font_body_fallback']
                        p.font.bold = (j == 0)

        if source:
            self._add_text(slide, f"Source: {source}",
                            Inches(0.5), Inches(6.7), Inches(12), Inches(0.3),
                            font_size=7, color=RGBColor(0x80, 0x80, 0x80),
                            font_name=self.theme['font_body_fallback'])

    def add_disclaimer_slide(self, text: str = '', title: str = 'Important Disclosures'):
        """免责声明尾页 — 小字体法律文本"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)

        self._add_shape(slide, MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0), self.slide_w, Inches(0.06),
                        self.theme['primary'])

        self._add_text(slide, title,
                        Inches(0.8), Inches(0.4), Inches(11), Inches(0.5),
                        font_size=18, bold=True, color=self.theme['primary'])

        default_text = (
            "This document is confidential and has been prepared by the Investment Banking Division solely for "
            "informational purposes. It is not intended as an offer or solicitation for the purchase or sale of "
            "any financial instrument. The information contained herein has been obtained from sources believed "
            "to be reliable but is not guaranteed as to its accuracy or completeness.\n\n"
            "Past performance is not indicative of future results. Any projections, estimates, forecasts, targets, "
            "prospects and/or opinions expressed in these materials are subject to change without notice and may "
            "differ or be contrary to opinions expressed by others.\n\n"
            "This presentation may not be reproduced, distributed or transmitted, in whole or in part, without "
            "the prior written consent of the Investment Banking Division. By accepting this document, the "
            "recipient agrees to be bound by the foregoing limitations."
        )

        self._add_text(slide, text or default_text,
                        Inches(0.8), Inches(1.2), Inches(11.5), Inches(5.5),
                        font_size=9, color=RGBColor(0x66, 0x66, 0x66),
                        font_name=self.theme['font_body_fallback'])

    def save(self, output_path: str):
        """保存 PPTX"""
        os.makedirs(os.path.dirname(output_path) or '.', exist_ok=True)
        self.prs.save(output_path)
        print(f"✅ PPTX saved: {output_path} ({os.path.getsize(output_path) / 1024:.0f} KB)")


# ═══════════════════════════════════════════
# 渲染入口
# ═══════════════════════════════════════════

LAYOUT_MAP = {
    'title': 'add_title_slide',
    'section': 'add_section_slide',
    'content': 'add_content_slide',
    'two_column': 'add_two_column_slide',
    'table': 'add_table_slide',
    'summary': 'add_summary_slide',
}


def render_pptx(data: Dict, output_path: str, template: str = 'default'):
    """
    根据 JSON 数据渲染 PPTX。

    data = {
        "title": "报告标题",
        "subtitle": "副标题",
        "author": "作者",
        "date": "日期",
        "theme": "cicc",  # 可选，覆盖 template 参数
        "slides": [
            {"layout": "title", "title": "...", ...},
            {"layout": "content", "title": "...", "content": "...", "bullets": [...]},
            {"layout": "table", "title": "...", "headers": [...], "rows": [...]},
        ]
    }
    """
    theme = data.get('theme', template)
    builder = PptxBuilder(theme_name=theme)

    # 封面
    if data.get('title'):
        builder.add_title_slide(
            title=data.get('title', ''),
            subtitle=data.get('subtitle', ''),
            author=data.get('author', ''),
            date=data.get('date', ''),
        )

    # 各页
    for slide_data in data.get('slides', []):
        layout = slide_data.get('layout', 'content')

        if layout == 'title':
            builder.add_title_slide(
                title=slide_data.get('title', ''),
                subtitle=slide_data.get('subtitle', ''),
            )
        elif layout == 'section':
            builder.add_section_slide(
                title=slide_data.get('title', ''),
                subtitle=slide_data.get('subtitle', ''),
            )
        elif layout == 'content':
            builder.add_content_slide(
                title=slide_data.get('title', ''),
                content=slide_data.get('content', ''),
                bullets=slide_data.get('bullets'),
                image=slide_data.get('image'),
                notes=slide_data.get('notes'),
            )
        elif layout == 'two_column':
            builder.add_two_column_slide(
                title=slide_data.get('title', ''),
                left_content=slide_data.get('left_content', slide_data.get('content', '')),
                right_content=slide_data.get('right_content', ''),
                right_image=slide_data.get('image'),
            )
        elif layout == 'table':
            builder.add_table_slide(
                title=slide_data.get('title', ''),
                headers=slide_data.get('headers', []),
                rows=slide_data.get('rows', []),
            )
        elif layout == 'summary':
            builder.add_summary_slide(
                title=slide_data.get('title', ''),
                points=slide_data.get('points', slide_data.get('bullets', [])),
            )

        elif layout == 'chart':
            builder.add_chart_slide(
                title=slide_data.get('title', ''),
                image=slide_data.get('image', ''),
                caption=slide_data.get('caption', ''),
            )
        elif layout == 'kpi':
            builder.add_kpi_slide(
                title=slide_data.get('title', ''),
                kpis=slide_data.get('kpis', []),
            )
        elif layout == 'comparison':
            builder.add_comparison_slide(
                title=slide_data.get('title', ''),
                left_title=slide_data.get('left_title', ''),
                left_items=slide_data.get('left_items', []),
                right_title=slide_data.get('right_title', ''),
                right_items=slide_data.get('right_items', []),
            )
        elif layout == 'timeline':
            builder.add_timeline_slide(
                title=slide_data.get('title', ''),
                events=slide_data.get('events', []),
            )
        elif layout == 'quote':
            builder.add_quote_slide(
                quote=slide_data.get('quote', slide_data.get('content', '')),
                author=slide_data.get('author', ''),
                source=slide_data.get('source', ''),
            )
        elif layout == 'end':
            builder.add_end_slide(
                title=slide_data.get('title', 'Thank You'),
                subtitle=slide_data.get('subtitle', ''),
                contact=slide_data.get('contact', ''),
            )

        # 投行 Pitch Book 专用
        elif layout == 'comparable_companies':
            builder.add_comparable_companies_slide(
                title=slide_data.get('title', 'Comparable Companies Analysis'),
                headers=slide_data.get('headers', []),
                rows=slide_data.get('rows', []),
                summary_rows=slide_data.get('summary_rows'),
                source=slide_data.get('source', ''),
            )
        elif layout == 'football_field':
            builder.add_football_field_slide(
                title=slide_data.get('title', 'Valuation Summary'),
                ranges=slide_data.get('ranges', []),
                current_price=slide_data.get('current_price'),
                currency=slide_data.get('currency', '$'),
                source=slide_data.get('source', ''),
            )
        elif layout == 'sources_uses':
            builder.add_sources_uses_slide(
                title=slide_data.get('title', 'Sources & Uses'),
                sources=slide_data.get('sources', []),
                uses=slide_data.get('uses', []),
                currency=slide_data.get('currency', '$m'),
                source_note=slide_data.get('source', ''),
            )

        elif layout == 'sensitivity_matrix':
            builder.add_sensitivity_matrix_slide(
                title=slide_data.get('title', 'Sensitivity Analysis'),
                row_label=slide_data.get('row_label', 'Terminal Growth Rate'),
                col_label=slide_data.get('col_label', 'WACC'),
                row_values=slide_data.get('row_values', []),
                col_values=slide_data.get('col_values', []),
                matrix=slide_data.get('matrix', []),
                highlight_row=slide_data.get('highlight_row'),
                highlight_col=slide_data.get('highlight_col'),
                source=slide_data.get('source', ''),
            )
        elif layout == 'transaction_overview':
            builder.add_transaction_overview_slide(
                title=slide_data.get('title', 'Transaction Overview'),
                key_points=slide_data.get('key_points', slide_data.get('bullets', [])),
                terms=slide_data.get('terms', []),
                source=slide_data.get('source', ''),
            )
        elif layout == 'disclaimer':
            builder.add_disclaimer_slide(
                text=slide_data.get('content', slide_data.get('text', '')),
                title=slide_data.get('title', 'Important Disclosures'),
            )

    builder.save(output_path)
    return output_path
