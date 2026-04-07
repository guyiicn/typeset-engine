#!/usr/bin/env python3
"""
文件对比工具 — 支持 PDF/PPTX/DOCX 的内容和视觉对比。

模式:
  text   — 提取文本内容对比（快速，精确定位文字变化）
  visual — 转为图片逐页对比（直观，看排版/字体差异）
  both   — 两者都做
"""

import os
import subprocess
import tempfile
from pathlib import Path
from typing import Optional


def _extract_text_pdf(filepath: str) -> str:
    """从 PDF 提取文本"""
    try:
        import pdfplumber
        texts = []
        with pdfplumber.open(filepath) as pdf:
            for i, page in enumerate(pdf.pages):
                t = page.extract_text()
                if t:
                    texts.append(f"--- Page {i+1} ---\n{t}")
        return '\n\n'.join(texts)
    except Exception as e:
        return f"[Text extraction failed: {e}]"


def _extract_text_pptx(filepath: str) -> str:
    """从 PPTX 提取文本"""
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        texts = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    slide_text.append(shape.text)
            if slide_text:
                texts.append(f"--- Slide {i+1} ---\n" + '\n'.join(slide_text))
        return '\n\n'.join(texts)
    except Exception as e:
        return f"[Text extraction failed: {e}]"


def _extract_text_docx(filepath: str) -> str:
    """从 DOCX 提取文本"""
    try:
        from docx import Document
        doc = Document(filepath)
        texts = [p.text for p in doc.paragraphs if p.text.strip()]
        return '\n'.join(texts)
    except Exception as e:
        return f"[Text extraction failed: {e}]"


def _extract_text(filepath: str) -> str:
    ext = Path(filepath).suffix.lower()
    if ext == '.pdf':
        return _extract_text_pdf(filepath)
    elif ext == '.pptx':
        return _extract_text_pptx(filepath)
    elif ext == '.docx':
        return _extract_text_docx(filepath)
    else:
        with open(filepath) as f:
            return f.read()


def _text_diff(old_text: str, new_text: str) -> str:
    """生成文本差异"""
    import difflib
    old_lines = old_text.splitlines()
    new_lines = new_text.splitlines()

    diff = difflib.unified_diff(old_lines, new_lines,
                                 fromfile='old', tofile='new',
                                 lineterm='')
    diff_text = '\n'.join(diff)

    if not diff_text:
        return "✅ No differences found (identical content)"

    # 统计
    added = sum(1 for l in diff_text.split('\n') if l.startswith('+') and not l.startswith('+++'))
    removed = sum(1 for l in diff_text.split('\n') if l.startswith('-') and not l.startswith('---'))

    header = f"📝 Text Diff: +{added} lines added, -{removed} lines removed\n{'='*60}\n"
    return header + diff_text


def _visual_diff(old_file: str, new_file: str, output: str) -> str:
    """视觉对比 — PDF 转图片后逐页 diff"""
    try:
        from PIL import Image
        import io

        tmpdir = tempfile.mkdtemp()

        # PDF → 图片 (用 pdftoppm)
        def pdf_to_images(pdf_path, prefix):
            subprocess.run(
                ['pdftoppm', '-png', '-r', '150', pdf_path, f'{tmpdir}/{prefix}'],
                capture_output=True, timeout=30
            )
            imgs = sorted(Path(tmpdir).glob(f'{prefix}-*.png'))
            return [str(p) for p in imgs]

        ext = Path(old_file).suffix.lower()
        if ext == '.pdf':
            old_imgs = pdf_to_images(old_file, 'old')
            new_imgs = pdf_to_images(new_file, 'new')
        else:
            return "⚠️ Visual diff currently supports PDF only"

        if not old_imgs or not new_imgs:
            return "⚠️ Could not convert PDF to images"

        # 逐页对比
        results = []
        max_pages = max(len(old_imgs), len(new_imgs))

        for i in range(max_pages):
            if i < len(old_imgs) and i < len(new_imgs):
                # 用 ImageMagick compare
                diff_img = f'{tmpdir}/diff_page{i+1}.png'
                result = subprocess.run(
                    ['compare', '-metric', 'AE', old_imgs[i], new_imgs[i], diff_img],
                    capture_output=True, text=True, timeout=30
                )
                pixel_diff = result.stderr.strip()
                results.append(f"Page {i+1}: {pixel_diff} pixels differ")
            elif i < len(new_imgs):
                results.append(f"Page {i+1}: NEW PAGE (not in old)")
            else:
                results.append(f"Page {i+1}: REMOVED (not in new)")

        # 合并 diff 图片
        diff_imgs = sorted(Path(tmpdir).glob('diff_page*.png'))
        if diff_imgs and output:
            # 拼接所有 diff 页
            images = [Image.open(p) for p in diff_imgs]
            if images:
                total_h = sum(img.height for img in images)
                max_w = max(img.width for img in images)
                combined = Image.new('RGB', (max_w, total_h), 'white')
                y = 0
                for img in images:
                    combined.paste(img, (0, y))
                    y += img.height
                combined.save(output)
                results.append(f"\n📸 Visual diff saved: {output}")

        # 清理
        import shutil
        shutil.rmtree(tmpdir, ignore_errors=True)

        header = f"🔍 Visual Diff: {max_pages} pages compared\n{'='*60}\n"
        return header + '\n'.join(results)

    except Exception as e:
        return f"⚠️ Visual diff failed: {e}"


def compare_files(old_file: str, new_file: str,
                   output: Optional[str] = None, mode: str = 'both') -> str:
    """
    对比两个文档。

    Args:
        old_file: 旧版文件路径
        new_file: 新版文件路径
        output: 差异报告输出路径（用于视觉对比图）
        mode: 'text' / 'visual' / 'both'

    Returns:
        差异报告文本
    """
    results = []
    results.append(f"📄 Comparing: {Path(old_file).name} → {Path(new_file).name}")
    results.append(f"Mode: {mode}\n")

    if mode in ('text', 'both'):
        old_text = _extract_text(old_file)
        new_text = _extract_text(new_file)
        results.append(_text_diff(old_text, new_text))

    if mode in ('visual', 'both'):
        out = output or f'/tmp/diff_{Path(new_file).stem}.png'
        results.append('\n' + _visual_diff(old_file, new_file, out))

    return '\n'.join(results)
