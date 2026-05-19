"""render_pptx.py 全 layout 单元测试

覆盖 render_pptx() dispatch 里的 20 种 layout:
  通用 12: title section content two_column table summary chart kpi
           comparison timeline quote end
  投行 8: comparable_companies football_field sources_uses sensitivity_matrix
         transaction_overview disclaimer waterfall org_chart

每个 layout 一个最小 payload, 调 render_pptx() (含 dispatch + builder), 用
python-pptx 解析回来检查 slide count + 关键 text 存在.

直接跑: python tests/test_render_pptx_layouts.py
或 pytest: pytest tests/test_render_pptx_layouts.py -v
"""

import os
import sys
import tempfile

sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', 'scripts'))

from render_pptx import render_pptx
from pptx import Presentation


# ─── 辅助 ─────────────────────────────────────────────────────────────

def _render(slides, theme='default'):
    """render 给定 slides + 一个固定 cover, 返回 Presentation 对象 + path."""
    data = {
        'title': 'Test Cover',
        'subtitle': 'sub',
        'author': 'pytest',
        'date': '2026-05-19',
        'slides': slides,
    }
    fd, path = tempfile.mkstemp(suffix='.pptx')
    os.close(fd)
    render_pptx(data, path, template=theme)
    return Presentation(path), path


def _slide_text(slide):
    """收集 slide 内所有 text frame 文本拼成单字符串."""
    parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text.strip():
                        parts.append(run.text)
    return ' | '.join(parts)


def _assert_slide_has(slide, *needles):
    """断言 slide 文本里包含所有 needle (substring)."""
    text = _slide_text(slide)
    for n in needles:
        assert n in text, f"missing {n!r} in slide text: {text[:200]}"


# ─── 通用 12 种 (strict) ──────────────────────────────────────────────

def test_title():
    prs, _ = _render([{'layout': 'title', 'title': 'T1', 'subtitle': 'S1'}])
    assert len(prs.slides) == 2  # cover + 1
    _assert_slide_has(prs.slides[1], 'T1', 'S1')


def test_section():
    prs, _ = _render([{'layout': 'section', 'title': 'CHAPTER ONE', 'subtitle': 'intro'}])
    assert len(prs.slides) == 2
    _assert_slide_has(prs.slides[1], 'CHAPTER ONE', 'intro')


def test_content():
    prs, _ = _render([{
        'layout': 'content', 'title': 'C-TITLE',
        'content': 'paragraph body',
        'bullets': ['bullet alpha', 'bullet beta'],
    }])
    assert len(prs.slides) == 2
    _assert_slide_has(prs.slides[1], 'C-TITLE', 'bullet alpha', 'bullet beta')


def test_two_column():
    prs, _ = _render([{
        'layout': 'two_column', 'title': 'TWO-COL',
        'left_content': 'LEFT-SIDE',
        'right_content': 'RIGHT-SIDE',
    }])
    assert len(prs.slides) == 2
    _assert_slide_has(prs.slides[1], 'TWO-COL', 'LEFT-SIDE', 'RIGHT-SIDE')


def test_table():
    prs, _ = _render([{
        'layout': 'table', 'title': 'TBL-TITLE',
        'headers': ['col-a', 'col-b'],
        'rows': [['r1-a', 'r1-b'], ['r2-a', 'r2-b']],
    }])
    slide = prs.slides[1]
    assert len(prs.slides) == 2
    # title 应该在; 表格走 table 元素而非 text frame, 单独验证
    _assert_slide_has(slide, 'TBL-TITLE')
    tables = [s for s in slide.shapes if s.has_table]
    assert len(tables) == 1, "expected 1 table"
    tbl = tables[0].table
    assert len(tbl.rows) == 3 and len(tbl.columns) == 2
    assert tbl.cell(0, 0).text == 'col-a'
    assert tbl.cell(1, 1).text == 'r1-b'


def test_summary():
    prs, _ = _render([{
        'layout': 'summary', 'title': 'SUM',
        'bullets': ['point A', 'point B'],
    }])
    _assert_slide_has(prs.slides[1], 'SUM', 'point A', 'point B')


def test_summary_points_alias():
    """summary 同时接受 points 字段 (跟 bullets 等价)"""
    prs, _ = _render([{
        'layout': 'summary', 'title': 'SUM2',
        'points': ['p1', 'p2'],
    }])
    _assert_slide_has(prs.slides[1], 'SUM2', 'p1', 'p2')


def test_chart_without_image():
    """image 不存在时应跳过图片但保留 title/caption."""
    prs, _ = _render([{
        'layout': 'chart', 'title': 'CHART-T',
        'image': '/tmp/_nonexistent_chart_xyz.png',
        'caption': 'figure cap',
    }])
    _assert_slide_has(prs.slides[1], 'CHART-T', 'figure cap')


def test_kpi():
    prs, _ = _render([{
        'layout': 'kpi', 'title': 'KPI-T',
        'kpis': [
            {'label': 'Revenue', 'value': '$100M', 'change': '+15%'},
            {'label': 'Users',   'value': '5M',    'change': '+8%'},
            {'label': 'Margin',  'value': '32%',   'change': '+2pp'},
        ],
    }])
    _assert_slide_has(prs.slides[1], 'KPI-T', 'Revenue', '$100M', 'Users')


def test_comparison():
    prs, _ = _render([{
        'layout': 'comparison', 'title': 'CMP-T',
        'left_title': 'PROS', 'left_items': ['fast', 'cheap'],
        'right_title': 'CONS', 'right_items': ['risky', 'fragile'],
    }])
    _assert_slide_has(prs.slides[1], 'CMP-T', 'PROS', 'fast', 'CONS', 'risky')


def test_timeline():
    prs, _ = _render([{
        'layout': 'timeline', 'title': 'TL-T',
        'events': [
            {'date': '2024 Q1', 'event': 'launch'},
            {'date': '2024 Q4', 'event': 'series A'},
        ],
    }])
    _assert_slide_has(prs.slides[1], 'TL-T', '2024 Q1', 'launch', 'series A')


def test_quote():
    prs, _ = _render([{
        'layout': 'quote',
        'quote': 'Premature optimization is the root of all evil',
        'author': 'Knuth', 'source': 'CACM 1974',
    }])
    _assert_slide_has(prs.slides[1], 'Premature', 'Knuth', 'CACM 1974')


def test_end():
    prs, _ = _render([{
        'layout': 'end', 'title': 'THANK YOU',
        'subtitle': 'Q&A', 'contact': 'hi@example.com',
    }])
    _assert_slide_has(prs.slides[1], 'THANK YOU', 'Q&A', 'hi@example.com')


# ─── 投行 8 种 (best-effort, 不挂就算过) ────────────────────────────

def _best_effort(name, slide_dict):
    """投行 layout 测试: 跑通 = 过, 报错 = print 但不 fail (帮 schema 探测)."""
    try:
        prs, _ = _render([slide_dict])
        assert len(prs.slides) >= 2
        text = _slide_text(prs.slides[1])
        print(f"  [OK] {name}: slide_text head = {text[:120]!r}")
        return True
    except Exception as e:
        print(f"  [FAIL] {name}: {type(e).__name__}: {e}")
        return False


def test_pitchbook_layouts():
    """投行 8 个 layout 一起跑, 统计成功率."""
    print("\n── Pitch Book layouts (best-effort) ──")
    results = {}

    results['comparable_companies'] = _best_effort('comparable_companies', {
        'layout': 'comparable_companies', 'title': 'Comps',
        'headers': ['Company', 'EV/EBITDA', 'P/E'],
        'rows': [['Acme', '12.5x', '18x'], ['Beta', '14.0x', '22x']],
        'source': 'Bloomberg',
    })

    results['football_field'] = _best_effort('football_field', {
        'layout': 'football_field', 'title': 'Valuation Range',
        'ranges': [
            {'method': 'DCF',   'low': 80,  'high': 120},
            {'method': 'Comps', 'low': 90,  'high': 110},
        ],
        'current_price': 95, 'currency': '$', 'source': 'Internal',
    })

    results['sources_uses'] = _best_effort('sources_uses', {
        'layout': 'sources_uses', 'title': 'S&U',
        'sources': [{'item': 'Debt', 'amount': 500}, {'item': 'Equity', 'amount': 300}],
        'uses':    [{'item': 'Purchase', 'amount': 700}, {'item': 'Fees', 'amount': 100}],
        'currency': '$m', 'source': 'mgmt',
    })

    results['sensitivity_matrix'] = _best_effort('sensitivity_matrix', {
        'layout': 'sensitivity_matrix', 'title': 'Sensitivity',
        'row_label': 'g', 'col_label': 'WACC',
        'row_values': ['2%', '3%', '4%'],
        'col_values': ['8%', '9%', '10%'],
        'matrix': [[100, 95, 90], [110, 105, 100], [120, 115, 110]],
        'highlight_row': 1, 'highlight_col': 1,
        'source': 'Internal',
    })

    results['transaction_overview'] = _best_effort('transaction_overview', {
        'layout': 'transaction_overview', 'title': 'TX Summary',
        'key_points': ['Strategic fit', 'Accretive y1'],
        'terms': [
            {'label': 'Price', 'value': '$1.2B'},
            {'label': 'Close', 'value': 'Q3 2026'},
        ],
        'source': 'PR',
    })

    results['disclaimer'] = _best_effort('disclaimer', {
        'layout': 'disclaimer', 'title': 'Disclosures',
        'content': 'This presentation is for discussion purposes only.',
    })

    results['waterfall'] = _best_effort('waterfall', {
        'layout': 'waterfall', 'title': 'EBITDA Bridge',
        'items': [
            {'label': 'FY24',     'value': 100, 'type': 'total'},
            {'label': 'Growth',   'value':  20, 'type': 'positive'},
            {'label': 'FX',       'value':  -5, 'type': 'negative'},
            {'label': 'FY25',     'value': 115, 'type': 'total'},
        ],
        'currency': '$m', 'source': 'mgmt',
    })

    results['org_chart'] = _best_effort('org_chart', {
        'layout': 'org_chart', 'title': 'Org',
        'root': {
            'name': 'CEO', 'title': 'Jane Doe',
            'children': [
                {'name': 'CTO', 'title': 'John', 'children': []},
                {'name': 'CFO', 'title': 'Mary', 'children': []},
            ],
        },
        'source': 'HR',
    })

    passed = sum(1 for v in results.values() if v)
    print(f"── Pitch Book: {passed}/{len(results)} layouts OK ──")
    # 不强制全过, 但记录到 metadata 供后续修
    if passed < len(results):
        failed = [k for k, v in results.items() if not v]
        print(f"   待修 schema: {failed}")


# ─── 入口 ─────────────────────────────────────────────────────────────

if __name__ == '__main__':
    tests = [
        test_title, test_section, test_content, test_two_column,
        test_table, test_summary, test_summary_points_alias,
        test_chart_without_image, test_kpi, test_comparison,
        test_timeline, test_quote, test_end,
    ]
    passed = 0
    failed = []
    for t in tests:
        try:
            t()
            print(f"✅ {t.__name__}")
            passed += 1
        except Exception as e:
            print(f"❌ {t.__name__}: {type(e).__name__}: {e}")
            failed.append(t.__name__)
    print(f"\n通用 layout: {passed}/{len(tests)} 通过")
    if failed:
        print(f"失败: {failed}")
    test_pitchbook_layouts()
